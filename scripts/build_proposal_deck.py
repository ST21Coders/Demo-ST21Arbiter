"""Build the ARBITER client proposal deck.

Two slides:
  1. 2x2 quadrant — Problems / ARBITER Solution / Use Cases / Add-ons
  2. Simplified 3-tier architecture — User & Edge / Compute & Agents / Data & AI

Visual theme mirrors Documents/ST21_Offerings.pptx (SMARTEK21 corporate template):
  Fonts:   Aptos Display (headings), Aptos (body)
  Size:    16:9 widescreen (13.333" x 7.5")
  Palette: #0E2841 navy, #156082 steel, #E97132 orange, #196B24 green,
           #E8E8E8 grey, #FFFFFF white

Run from the project root with the project venv activated:
    python3 scripts/build_proposal_deck.py
"""
from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ──────────────────────────── theme ─────────────────────────────
NAVY   = RGBColor(0x0E, 0x28, 0x41)
STEEL  = RGBColor(0x15, 0x60, 0x82)
ORANGE = RGBColor(0xE9, 0x71, 0x32)
GREEN  = RGBColor(0x19, 0x6B, 0x24)
GREY   = RGBColor(0xE8, 0xE8, 0xE8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x20, 0x20, 0x20)

AWS = {
    "lambda":    RGBColor(0xED, 0x71, 0x00),
    "agentcore": RGBColor(0x01, 0xA8, 0x8D),
    "purple":    RGBColor(0x8C, 0x4F, 0xFF),
    "apigw":     RGBColor(0xE7, 0x15, 0x7B),
    "cognito":   RGBColor(0xDD, 0x34, 0x4C),
    "s3":        RGBColor(0x7A, 0xA1, 0x16),
    "ddb":       RGBColor(0xC9, 0x25, 0xD1),
    "user":      RGBColor(0x23, 0x2F, 0x3E),
}

HEAD_FONT = "Aptos Display"
BODY_FONT = "Aptos"

OUT = Path(__file__).resolve().parent.parent / "Documents" / "ARBITER_Client_Proposal.pptx"


# ──────────────────────────── primitives ────────────────────────
def fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape) -> None:
    shape.line.fill.background()


def set_text(
    frame,
    text: str,
    *,
    font: str = BODY_FONT,
    size: int = 11,
    color: RGBColor = DARK_TEXT,
    bold: bool = False,
    italic: bool = False,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
) -> None:
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = anchor
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_box(slide, left, top, width, height, color: RGBColor, *, rounded: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    fill(s, color)
    no_line(s)
    return s


def add_text(
    slide, left, top, width, height, text: str, **text_kwargs,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    set_text(tb.text_frame, text, **text_kwargs)
    return tb


def add_bullets(slide, left, top, width, height, bullets: list[str], *, size: int = 11) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.06)
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"•  {line}"
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = DARK_TEXT
        p.space_after = Pt(3)


def add_right_arrow(slide, left, top, width, height, color: RGBColor = STEEL):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    fill(a, color)
    no_line(a)
    return a


def add_tile(slide, left, top, width, height, color: RGBColor, label: str, sublabel: str = ""):
    """AWS service tile — rounded rectangle, white label, optional sublabel."""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill(rect, color)
    no_line(rect)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = HEAD_FONT
    r.font.size = Pt(11)
    r.font.color.rgb = WHITE
    r.font.bold = True
    if sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sublabel
        r2.font.name = BODY_FONT
        r2.font.size = Pt(8)
        r2.font.color.rgb = WHITE
    return rect


# ──────────────────────────── slide-level helpers ───────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title_bar(slide, title: str, subtitle: str) -> None:
    add_box(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), NAVY)
    add_text(
        slide, Inches(0.4), Inches(0.10), Inches(12.5), Inches(0.65), title,
        font=HEAD_FONT, size=24, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.40), subtitle,
        font=BODY_FONT, size=12, color=STEEL, italic=True, anchor=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide) -> None:
    today = date.today().isoformat()
    add_text(
        slide, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.3),
        f"ARBITER ST21  ·  Smartek21  ·  {today}",
        font=BODY_FONT, size=9, color=NAVY, italic=True,
    )


def add_quadrant(
    slide, left, top, width, height,
    header_color: RGBColor, header_text: str, bullets: list[str],
) -> None:
    # Outer panel (grey body)
    add_box(slide, left, top, width, height, GREY)
    # Header strip (colored)
    add_box(slide, left, top, width, Inches(0.45), header_color)
    # Header text
    add_text(
        slide, left + Inches(0.12), top, width - Inches(0.24), Inches(0.45),
        header_text, font=HEAD_FONT, size=15, color=WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Bullets — sit in the remaining body
    add_bullets(
        slide,
        left + Inches(0.08), top + Inches(0.48),
        width - Inches(0.16), height - Inches(0.52),
        bullets, size=11,
    )


def add_band_header(slide, left, top, width, color: RGBColor, label: str) -> None:
    add_box(slide, left, top, width, Inches(0.32), color)
    add_text(
        slide, left + Inches(0.12), top, width - Inches(0.24), Inches(0.32),
        label, font=HEAD_FONT, size=12, color=WHITE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ──────────────────────────── content ───────────────────────────
TITLE = "ARBITER — Smart AI Governance for Enterprise Policy"
SUBTITLE = (
    "Multi-agent reconciliation across SharePoint, AWS Config, and Zscaler "
    "— grounded, audited, persona-aware."
)

PROBLEMS = [
    "Policies live in 3+ silos — SharePoint docs, AWS Config rules, Zscaler ZIA allowlists. No single source of truth.",
    "Manual reconciliation across systems is slow, error-prone, and audit-vulnerable.",
    "Drift surfaces reactively during incidents or audits — never proactively.",
    "Compliance frameworks (PCI-DSS, NAIC, SOC 2) need traceable policy mapping that doesn't exist today.",
    "New analyst questions take days — \"is X allowed?\" requires cross-team email chains.",
]

SOLUTION = [
    "Multi-agent orchestration — one query fans out to SharePoint + AWS Config + Zscaler specialists in parallel.",
    "Grounded answers — Bedrock Knowledge Base cites the exact policy document and rule for every finding.",
    "Guardrails — Bedrock content / PII / denied-topic filters wrap every model call.",
    "Persona-aware access — CISO, SOC, GRC, Employee each see role-appropriate data via Cognito groups.",
    "Conversation memory — AgentCore Memory keeps 90-day session history; follow-ups remember prior turns.",
    "Foundation-model agnostic — Nova 2 Lite default; swap to Claude or other Bedrock models per workload.",
]

USECASES = [
    "CISO (Diana) — review critical findings, approve high-severity change requests, configure LLM controls.",
    "SOC Analyst (Marcus) — alert-driven dashboard + heat map, kick off remediation workflows.",
    "GRC Analyst (Priya) — compliance posture across PCI-DSS / NAIC / SOC 2; one-click audit export.",
    "Employee (Sarah) — chatbot self-service: \"why is github.com blocked for me?\" with cited policy answer.",
]

ADDONS = [
    "ITSM integration — ServiceNow / Jira change-request sync (auto-create CR from a finding).",
    "Conversational channels — Slack / Microsoft Teams chatbot using the same master orchestrator.",
    "Compliance frameworks — pre-mapped controls for PCI-DSS, HIPAA, SOC 2, NAIC, ISO 27001.",
    "New data sources — Okta, CrowdStrike, Splunk, Wiz, Prisma as additional specialists.",
    "Multi-level approval workflows — escalation chains for severity tiers.",
    "Cost & usage analytics — token / Bedrock spend by persona, top-queried policies dashboard.",
]


# ──────────────────────────── Slide 1 ───────────────────────────
def build_slide_1(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_title_bar(s, TITLE, SUBTITLE)

    # Quadrant grid — 13.333" wide, available 12.93" between L/R margins
    LEFT_MARGIN = Inches(0.2)
    GAP = Inches(0.2)
    TOP = Inches(1.45)
    W = Inches(6.46)   # (13.333 - 0.4 margins - 0.2 gap) / 2
    H = Inches(2.7)

    add_quadrant(s, LEFT_MARGIN,                  TOP,            W, H, STEEL,  "Policy Management Problems",     PROBLEMS)
    add_quadrant(s, LEFT_MARGIN + W + GAP,        TOP,            W, H, ORANGE, "ARBITER — Smart AI Governance",  SOLUTION)
    add_quadrant(s, LEFT_MARGIN,                  TOP + H + GAP,  W, H, STEEL,  "End-User Use Cases",             USECASES)
    add_quadrant(s, LEFT_MARGIN + W + GAP,        TOP + H + GAP,  W, H, GREEN,  "Possible Add-On Features",       ADDONS)

    add_footer(s)


# ──────────────────────────── Slide 2 ───────────────────────────
TITLE_2 = "ARBITER Solution Architecture"
SUBTITLE_2 = (
    "Three planes — User & Edge → Compute & Agents → Data & AI — glued by "
    "Cognito identity and Bedrock guardrails."
)


def build_slide_2(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, TITLE_2, SUBTITLE_2)

    L = Inches(0.25)
    R = Inches(13.083)
    W = R - L

    # ─── Band 1: USER & EDGE ──────────────────────────────────
    band1_top = Inches(1.45)
    band1_h = Inches(1.45)
    add_box(s, L, band1_top, W, band1_h, GREY)
    add_band_header(s, L, band1_top, W, NAVY, "USER & EDGE")

    tile_w = Inches(1.85)
    tile_h = Inches(0.72)
    tile_y = band1_top + Inches(0.45)
    arrow_w = Inches(0.32)
    arrow_h = Inches(0.28)
    arrow_y = tile_y + (tile_h - arrow_h) / 2
    cur = L + Inches(0.25)
    tiles_1 = [
        (AWS["user"],    "User",            "browser"),
        (AWS["purple"],  "CloudFront",      "+ S3 (UI)"),
        (AWS["cognito"], "Cognito",         "Hosted UI · groups"),
        (AWS["apigw"],   "API Gateway",     "REST · JWT auth"),
        (AWS["lambda"],  "Function URL",    "/chat · long ops"),
    ]
    for i, (color, label, sub) in enumerate(tiles_1):
        add_tile(s, cur, tile_y, tile_w, tile_h, color, label, sub)
        cur += tile_w
        if i < len(tiles_1) - 1:
            add_right_arrow(s, cur + Inches(0.05), arrow_y, arrow_w, arrow_h, STEEL)
            cur += Inches(0.32 + 0.10)
    add_text(
        s, L + Inches(0.25), band1_top + band1_h - Inches(0.30),
        W - Inches(0.5), Inches(0.25),
        "Bearer JWT on every call  ·  API GW for sub-29 s ops  ·  Function URL for long agent fan-outs (up to 15 min).",
        font=BODY_FONT, size=10, color=NAVY, italic=True,
    )

    # ─── Band 2: COMPUTE & AGENTS ─────────────────────────────
    band2_top = Inches(2.98)
    band2_h = Inches(2.05)
    add_box(s, L, band2_top, W, band2_h, GREY)
    add_band_header(s, L, band2_top, W, NAVY, "COMPUTE & AGENTS  —  Master + 3 Specialists pattern")

    # api_handler on the left, vertically centered in the body region
    body2_top = band2_top + Inches(0.35)
    body2_h = band2_h - Inches(0.35)
    handler_w = Inches(1.85)
    handler_h = Inches(0.85)
    handler_y = body2_top + (body2_h - handler_h) / 2
    handler_x = L + Inches(0.35)
    add_tile(s, handler_x, handler_y, handler_w, handler_h, AWS["lambda"],
             "api_handler", "Lambda · routes")

    # Arrow to Master
    master_x = handler_x + handler_w + Inches(0.55)
    master_w = Inches(2.4)
    master_h = Inches(0.85)
    master_y = handler_y
    a1_x = handler_x + handler_w + Inches(0.05)
    a1_y = handler_y + (handler_h - Inches(0.28)) / 2
    add_right_arrow(s, a1_x, a1_y, Inches(0.45), Inches(0.28), STEEL)
    add_tile(s, master_x, master_y, master_w, master_h, AWS["agentcore"],
             "Master Orchestrator", "AgentCore Runtime · Nova 2 Lite")

    # Three specialists fanning out to the right
    spec_x = master_x + master_w + Inches(0.65)
    spec_w = Inches(2.4)
    spec_h = Inches(0.55)
    spec_gap = Inches(0.10)
    spec_total = 3 * spec_h + 2 * spec_gap
    spec_top = body2_top + (body2_h - spec_total) / 2
    specs = [
        ("SharePoint Specialist", "policy docs"),
        ("AWS Config Specialist", "cloud rules"),
        ("Zscaler Specialist",    "URL & DLP"),
    ]
    for i, (label, sub) in enumerate(specs):
        y = spec_top + i * (spec_h + spec_gap)
        add_tile(s, spec_x, y, spec_w, spec_h, AWS["agentcore"], label, sub)
        # Fan-out arrow from master right-edge to specialist left-edge
        ax1 = master_x + master_w + Inches(0.05)
        ay1 = y + (spec_h - Inches(0.22)) / 2
        add_right_arrow(s, ax1, ay1, Inches(0.50), Inches(0.22), STEEL)

    add_text(
        s, L + Inches(0.35), band2_top + band2_h - Inches(0.28),
        W - Inches(0.7), Inches(0.25),
        "One question fans out to 3 specialists in parallel · each grounds its answer in its own data domain.",
        font=BODY_FONT, size=10, color=NAVY, italic=True,
    )

    # ─── Band 3: DATA & AI PLANE ──────────────────────────────
    band3_top = Inches(5.10)
    band3_h = Inches(1.85)
    add_box(s, L, band3_top, W, band3_h, GREY)
    add_band_header(s, L, band3_top, W, NAVY, "DATA & AI PLANE")

    # Two rows of tiles
    row_h = Inches(0.55)
    row_gap = Inches(0.08)
    row1_y = band3_top + Inches(0.42)
    row2_y = row1_y + row_h + row_gap

    data_tiles_row1 = [
        (AWS["agentcore"], "Bedrock KB",          "policy vectors"),
        (AWS["purple"],    "OpenSearch",          "vector index"),
        (AWS["agentcore"], "Bedrock Guardrails",  "PII · denied"),
        (AWS["agentcore"], "Nova 2 Lite",         "Foundation Model"),
    ]
    data_tiles_row2 = [
        (AWS["agentcore"], "AgentCore Memory",    "90-day history"),
        (AWS["ddb"],       "DynamoDB · 4 tables", "conflicts · CR · audit · sessions"),
    ]

    def lay_row(tiles, y, count):
        total_w = W - Inches(0.5)
        gap = Inches(0.18)
        each = (total_w - gap * (count - 1)) / count
        x = L + Inches(0.25)
        for color, label, sub in tiles:
            add_tile(s, x, y, each, row_h, color, label, sub)
            x += each + gap

    lay_row(data_tiles_row1, row1_y, 4)
    lay_row(data_tiles_row2, row2_y, 2)

    add_text(
        s, L + Inches(0.25), band3_top + band3_h - Inches(0.26),
        W - Inches(0.5), Inches(0.25),
        "KB grounds answers in actual policy text  ·  Guardrails filter every call  ·  DDB holds session metadata  ·  Memory keeps 90-day conversation history.",
        font=BODY_FONT, size=10, color=NAVY, italic=True,
    )

    add_footer(s)


# ──────────────────────────── main ──────────────────────────────
def main() -> None:
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333"
    prs.slide_height = Emu(6858000)   # 7.5"
    # python-pptx leaves type="screen4x3" by default; set to widescreen so
    # the metadata matches the cx/cy dimensions.
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    prs.part._element.find(f"{ns}sldSz").set("type", "screen16x9")

    build_slide_1(prs)
    build_slide_2(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
